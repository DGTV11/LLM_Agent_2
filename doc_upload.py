import base64
from io import BytesIO

import fitz
from docx import Document
from docx.oxml.ns import qn
from pptx import Presentation

from llm import call_vlm
from prompts import SPR_PROMPT


def vlm_process_image(b64_image, img_type):
    return call_vlm(
        [
            {"role": "system", "content": SPR_PROMPT},
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": ""},
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:image/{img_type};base64,{b64_image}"
                        },
                    },
                ],
            },
        ]
    )


def extract_docx_cell_content(cell, doc):
    chunks = []
    for para in cell.paragraphs:
        for run in para.runs:
            if run.text:
                chunks.append(run.text)

            drawing_elems = run._element.findall(
                ".//w:drawing", namespaces=run._element.nsmap
            )
            for drawing in drawing_elems:
                namespaces = drawing.nsmap.copy() if drawing.nsmap else {}
                if "a" not in namespaces:
                    namespaces["a"] = (
                        "http://schemas.openxmlformats.org/drawingml/2006/main"
                    )

                blip = drawing.find(".//a:blip", namespaces=namespaces)
                if blip is not None:
                    embed_rid = blip.get(qn("r:embed"))
                    image_part = doc.part.related_parts[embed_rid]
                    blob = image_part.blob
                    b64_blob = base64.b64encode(blob).decode("utf-8")
                    content_type = image_part.content_type

                    try:
                        textified_image = vlm_process_image(
                            b64_blob, content_type.replace("image/", "")
                        )
                    except Exception:
                        textified_image = "IMAGE FAILED"

                    chunks.append(
                        "\n".join(
                            [
                                "===IMAGE START===",
                                textified_image,
                                "===IMAGE END===",
                            ]
                        )
                    )
    return "\n".join(chunks)


def extract_pptx_cell_content(cell):
    chunks = []
    for para in cell.text_frame.paragraphs:
        if para.text:
            chunks.append(para.text)
    # Optional: handle shapes/images inside cells if present
    for shape in getattr(cell, "shapes", []):
        if hasattr(shape, "image"):
            try:
                textified_image = vlm_process_image(
                    base64.b64encode(shape.image.blob).decode("utf-8"),
                    shape.image.content_type.replace("image/", ""),
                )
            except Exception:
                textified_image = "IMAGE FAILED"
            chunks.append(
                "\n".join(["===IMAGE START===", textified_image, "===IMAGE END==="])
            )
    return "\n".join(chunks)


def table_to_md_docx(table, doc):
    md = []
    headers = [extract_docx_cell_content(cell, doc) for cell in table.rows[0].cells]
    md.append("| " + " | ".join(headers) + " |")
    md.append("| " + " | ".join("---" for _ in headers) + " |")

    for row in table.rows[1:]:
        row_text = [extract_docx_cell_content(cell, doc) for cell in row.cells]
        md.append("| " + " | ".join(row_text) + " |")
    return "\n".join(md)


def table_to_md_pptx(table):
    md = []
    headers = [extract_pptx_cell_content(cell) for cell in table.rows[0].cells]
    md.append("| " + " | ".join(headers) + " |")
    md.append("| " + " | ".join("---" for _ in headers) + " |")

    for row in table.rows[1:]:
        row_text = [extract_pptx_cell_content(cell) for cell in row.cells]
        md.append("| " + " | ".join(row_text) + " |")
    return "\n".join(md)


def process_file(file: bytes, content_type: str):
    # Text
    if content_type.startswith("text"):
        return file.decode("utf-8")

    # Images
    elif content_type == "image/png":
        return vlm_process_image(base64.b64encode(file).decode("utf-8"), "png")
    elif content_type == "image/jpeg":
        return vlm_process_image(base64.b64encode(file).decode("utf-8"), "jpeg")

    # PPTX
    elif (
        content_type
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        slides = Presentation(BytesIO(file)).slides
        chunks = []

        for slide in slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    chunks.append(shape.text)
                if shape.has_table:
                    md_table = table_to_md_pptx(shape.table)
                    chunks.append(
                        "\n".join(["===TABLE START===", md_table, "===TABLE END==="])
                    )
                if hasattr(shape, "image"):
                    try:
                        textified_image = vlm_process_image(
                            base64.b64encode(shape.image.blob).decode("utf-8"),
                            shape.image.content_type.replace("image/", ""),
                        )
                    except Exception:
                        textified_image = "IMAGE FAILED"
                    chunks.append(
                        "\n".join(
                            ["===IMAGE START===", textified_image, "===IMAGE END==="]
                        )
                    )

        return "\n".join(chunks)

    # DOCX
    elif (
        content_type
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ):
        doc = Document(BytesIO(file))
        chunks = []

        for child in doc.element.body:
            if child.tag.endswith("p"):  # paragraph
                para = doc.paragraphs[doc.element.body.index(child)]
                chunk = []

                for run in para.runs:
                    if run.text:
                        chunk.append(run.text)

                    drawing_elems = run._element.findall(
                        ".//w:drawing", namespaces=run._element.nsmap
                    )
                    for drawing in drawing_elems:
                        namespaces = drawing.nsmap.copy() if drawing.nsmap else {}
                        if "a" not in namespaces:
                            namespaces["a"] = (
                                "http://schemas.openxmlformats.org/drawingml/2006/main"
                            )

                        blip = drawing.find(".//a:blip", namespaces=namespaces)
                        if blip is not None:
                            embed_rid = blip.get(qn("r:embed"))
                            image_part = doc.part.related_parts[embed_rid]
                            blob = image_part.blob
                            b64_blob = base64.b64encode(blob).decode("utf-8")
                            content_type_img = image_part.content_type

                            try:
                                textified_image = vlm_process_image(
                                    b64_blob, content_type_img.replace("image/", "")
                                )
                            except Exception:
                                textified_image = "IMAGE FAILED"

                            chunk.append(
                                "\n".join(
                                    [
                                        "===IMAGE START===",
                                        textified_image,
                                        "===IMAGE END===",
                                    ]
                                )
                            )

                chunks.append("\n".join(chunk))

            elif child.tag.endswith("tbl"):  # table
                table = doc.tables[doc.element.body.index(child)]
                md_table = table_to_md_docx(table, doc)
                chunks.append(
                    "\n".join(["===TABLE START===", md_table, "===TABLE END==="])
                )

        return "\n\n".join(chunks)

    # PDF
    elif content_type == "application/pdf":
        doc = fitz.open("pdf", file)
        chunks = []

        for page in doc:
            chunks.append(page.get_text("text"))

            for img_index, img in enumerate(page.get_images(full=True)):
                xref = img[0]
                base_image = doc.extract_image(xref)
                blob = base_image["image"]
                b64_blob = base64.b64encode(blob).decode("utf-8")
                ext = base_image["ext"]

                try:
                    textified_image = vlm_process_image(b64_blob, ext)
                except Exception:
                    textified_image = "IMAGE FAILED"

                chunks.append(
                    "\n".join(["===IMAGE START===", textified_image, "===IMAGE END==="])
                )

        return "\n".join(chunks)

    # Fallback
    try:
        return file.decode("utf-8")
    except Exception:
        raise ValueError("Invalid content_type")
